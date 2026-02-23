#!/usr/bin/env python3
"""
MegaNet Executive Presentation
Dark-blue theme, high-level visuals for a non-technical audience.
All graphics built with python-pptx shapes & charts (no external images).
"""
from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x0A, 0x14, 0x2A)   # deep navy
NAVY        = RGBColor(0x0F, 0x1E, 0x3C)
ACCENT      = RGBColor(0x00, 0x8C, 0xFF)   # electric blue
ACCENT2     = RGBColor(0x00, 0xD4, 0xFF)   # cyan
GREEN       = RGBColor(0x00, 0xC8, 0x6E)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
RED         = RGBColor(0xE8, 0x20, 0x40)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xB0, 0xC4, 0xDE)
MID_GREY    = RGBColor(0x4A, 0x5E, 0x7A)
YELLOW      = RGBColor(0xFF, 0xD7, 0x00)


# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
W = Cm(33.87)
H = Cm(19.05)


# ── Helpers ───────────────────────────────────────────────────────────────────
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, alpha=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def circle(slide, cx, cy, r, color, line_color=None, line_width=None):
    from pptx.util import Pt as Pt_
    shp = slide.shapes.add_shape(
        9,  # oval
        cx - r, cy - r, 2 * r, 2 * r
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color
        if line_width:
            shp.line.width = line_width
    else:
        shp.line.fill.background()
    return shp


def line(slide, x1, y1, x2, y2, color=ACCENT, width=Pt(2)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def txt(slide, text, l, t, w, h,
        size=Pt(14), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def multi_para(slide, lines, l, t, w, h, wrap=True):
    """lines: [(text, size, color, bold, align)]"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    for i, (text, size, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if text == "":
            p.space_after = Pt(4)
            continue
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def header_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, W, Cm(2.4), NAVY)
    rect(slide, 0, Cm(2.4), W, Cm(0.15), ACCENT)
    txt(slide, title, Cm(1), Cm(0.25), W - Cm(2), Cm(1.6),
        size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Cm(1), Cm(1.55), W - Cm(2), Cm(0.9),
            size=Pt(13), color=ACCENT2)


def footer_bar(slide, text="MegaNet MVP  |  Executive Briefing  |  2026"):
    rect(slide, 0, H - Cm(0.9), W, Cm(0.9), NAVY)
    txt(slide, text, Cm(1), H - Cm(0.85), W - Cm(2), Cm(0.8),
        size=Pt(10), color=LIGHT_GREY, align=PP_ALIGN.CENTER)


def pill(slide, text, l, t, w, h, bg_color, text_color=WHITE, size=Pt(13), bold=True):
    from pptx.util import Pt as Pt_
    shp = slide.shapes.add_shape(9, l, t, w, h)  # rounded via oval; use rect w/ rounding
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg_color
    shp.line.fill.background()
    # Add text via textbox on top
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    return shp, box


def rounded_rect(slide, l, t, w, h, color, radius=Cm(0.3)):
    """Approximate rounded rect using pptx rounded rectangle (shape type 5)."""
    shp = slide.shapes.add_shape(5, l, t, w, h)  # 5 = rounded rectangle
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    # Set corner radius via XML
    try:
        adj = shp.adjustments
        adj[0] = 0.1
    except Exception:
        pass
    return shp


def arrow_right(slide, x1, y, x2, color=ACCENT, width=Pt(2.5)):
    """Draw right-pointing arrow as connector."""
    conn = slide.shapes.add_connector(1, x1, y, x2, y)
    conn.line.color.rgb = color
    conn.line.width = width
    return conn


def icon_node(slide, cx, cy, label, color=ACCENT, r=Cm(0.6), lbl_color=WHITE, lbl_size=Pt(11)):
    """Draw a circle node with label below."""
    circle(slide, cx, cy, r, color)
    txt(slide, label,
        cx - Cm(1.2), cy + r + Cm(0.05),
        Cm(2.4), Cm(0.6),
        size=lbl_size, color=lbl_color, align=PP_ALIGN.CENTER)


def draw_wave(slide, cx, cy, amplitude=Cm(0.25), n_waves=3, color=ACCENT2, width=Pt(1.5)):
    """Simulate a radio wave with connected short line segments."""
    import math
    steps = n_waves * 8
    total_w = Cm(2.5)
    dx = total_w / steps
    pts = []
    for i in range(steps + 1):
        x = cx - total_w / 2 + i * dx
        y = cy + amplitude * math.sin(i / steps * n_waves * 2 * math.pi)
        pts.append((x, y))
    for i in range(len(pts) - 1):
        conn = slide.shapes.add_connector(1, pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1])
        conn.line.color.rgb = color
        conn.line.width = width


# ── Bar chart helper ──────────────────────────────────────────────────────────
def add_bar_chart(slide, l, t, w, h, categories, series_name, values, color=ACCENT):
    from pptx.util import Pt as Pt_
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, chart_data
    ).chart
    # Style
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    for series in plot.series:
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = color
    chart.category_axis.tick_labels.font.size = Pt_(10)
    chart.category_axis.tick_labels.font.color.rgb = WHITE
    chart.value_axis.tick_labels.font.size = Pt_(10)
    chart.value_axis.tick_labels.font.color.rgb = WHITE
    # Transparent background
    return chart


def add_donut_chart(slide, l, t, w, h, categories, values, colors_list):
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series("", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, l, t, w, h, chart_data
    ).chart
    chart.has_legend = True
    chart.has_title = False
    legend = chart.legend
    legend.font.size = Pt(11)
    legend.font.color.rgb = WHITE
    return chart


# =============================================================================
# SLIDES
# =============================================================================

def slide_cover(prs):
    s = new_slide(prs)
    bg(s, DARK_BG)

    # Gradient-like layered rectangles for visual depth
    rect(s, 0, 0, W, H, DARK_BG)
    rect(s, W * 0.55, 0, W * 0.45, H, RGBColor(0x05, 0x0E, 0x22))

    # Accent lines
    for i in range(6):
        c = RGBColor(0, max(0, 80 - i * 12), min(255, 140 + i * 18))
        rect(s, W * 0.55 + Cm(i * 0.9), 0, Cm(0.3), H, c)

    # Logo-like element (circle cluster)
    cx, cy = W * 0.77, H * 0.45
    for r_size, col in [
        (Cm(4.5), RGBColor(0, 30, 70)),
        (Cm(3.2), RGBColor(0, 60, 130)),
        (Cm(2.0), ACCENT),
        (Cm(0.9), ACCENT2),
    ]:
        circle(s, cx, cy, r_size, col)

    # Satellite nodes around centre
    import math
    for angle, nc in [(30, GREEN), (150, ORANGE), (270, YELLOW)]:
        nx = cx + Cm(3.5) * math.cos(math.radians(angle))
        ny = cy + Cm(3.5) * math.sin(math.radians(angle))
        line(s, cx, cy, nx, ny, ACCENT2, Pt(1.5))
        circle(s, nx, ny, Cm(0.4), nc)

    # Title
    multi_para(s, [
        ("MegaNet", Pt(60), WHITE, True, PP_ALIGN.LEFT),
        ("Internet Paralela Descentralizada", Pt(24), ACCENT2, False, PP_ALIGN.LEFT),
        ("", Pt(10), WHITE, False, PP_ALIGN.LEFT),
        ("Executive Briefing  |  2026", Pt(15), LIGHT_GREY, False, PP_ALIGN.LEFT),
    ], Cm(1.5), Cm(4.5), Cm(18), Cm(10))

    # Tagline box
    shp = rect(s, Cm(0), H - Cm(3.2), Cm(20), Cm(2.3), RGBColor(0, 100, 200))
    txt(s, "The internet reimagined. No IPs. No censorship. No single point of failure.",
        Cm(0.6), H - Cm(3.1), Cm(19), Cm(2.1),
        size=Pt(15), bold=True, color=WHITE, wrap=True)

    footer_bar(s, "CONFIDENTIAL  |  MegaNet MVP  |  Executive Briefing  |  2026")


def slide_problem(prs):
    s = new_slide(prs)
    bg(s)
    header_bar(s, "The Problem", "Why the current internet is broken")

    # Central hub (represents centralised internet)
    cx, cy = W / 2, H * 0.57
    circle(s, cx, cy, Cm(2.0), RGBColor(0x60, 0x10, 0x10))
    txt(s, "Centralised\nControl", cx - Cm(1.2), cy - Cm(0.6), Cm(2.4), Cm(1.2),
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Problems radiating out
    import math
    problems = [
        (30,  "IP Addresses\nassigned by\ngovernments",   RED),
        (90,  "DNS easily\nblocked &\nmonitored",         ORANGE),
        (150, "ISP can\ninspect all\ntraffic",            RED),
        (210, "BGP routing\nhijackable",                  ORANGE),
        (270, "Cloud platforms\ncan deplatform\nanyone",  RED),
        (330, "Single point\nof failure\n& censorship",   ORANGE),
    ]
    r_out = Cm(5.2)
    r_box_w = Cm(4.2)
    r_box_h = Cm(1.8)
    for angle, label, col in problems:
        rad = math.radians(angle)
        nx = cx + r_out * math.cos(rad)
        ny = cy + r_out * math.sin(rad)
        # line from centre
        lx = cx + Cm(2.1) * math.cos(rad)
        ly = cy + Cm(2.1) * math.sin(rad)
        line(s, lx, ly, nx - Cm(0.5) * math.cos(rad), ny - Cm(0.5) * math.sin(rad),
             col, Pt(1.5))
        circle(s, nx, ny, Cm(0.5), col)
        # label box
        bx = nx - r_box_w / 2
        by = ny + Cm(0.6)
        txt(s, label, bx, by, r_box_w, r_box_h,
            size=Pt(10.5), color=WHITE, align=PP_ALIGN.CENTER)

    # Big X over centre
    txt(s, "X", cx - Cm(0.5), cy - Cm(0.7), Cm(1), Cm(1.4),
        size=Pt(36), bold=True, color=RGBColor(255, 50, 50), align=PP_ALIGN.CENTER)

    footer_bar(s)


def slide_vision(prs):
    s = new_slide(prs)
    bg(s)
    header_bar(s, "The MegaNet Vision", "A parallel internet — no IPs, no DNS, no censors")

    # Two columns: OLD vs NEW
    mid_x = W / 2
    col_w = Cm(14.5)
    top_y = Cm(2.8)

    # OLD column header
    rect(s, Cm(0.5), top_y, col_w, Cm(0.9), RED)
    txt(s, "Traditional Internet", Cm(0.5), top_y, col_w, Cm(0.9),
        size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # NEW column header
    rect(s, mid_x + Cm(0.5), top_y, col_w, Cm(0.9), GREEN)
    txt(s, "MegaNet", mid_x + Cm(0.5), top_y, col_w, Cm(0.9),
        size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Comparison rows
    rows = [
        ("Identity",   "IP address (government/ISP assigned)",  "Public key (self-sovereign, cryptographic)"),
        ("Naming",     "DNS (centralised, blockable)",           "Blockchain registry (distributed, immutable)"),
        ("Transport",  "TCP/IP over ISP cables",                 "LoRa radio (unlicensed, no ISP needed)"),
        ("Privacy",    "Unencrypted by default",                 "E2E encrypted always (ChaCha20-Poly1305)"),
        ("Censorship", "Single points of control",               "No central entity can block nodes"),
        ("Data",       "Stored on centralised servers",          "'Homeless' data: distributed fragments"),
    ]
    row_h = Cm(1.55)
    for i, (aspect, old, new) in enumerate(rows):
        y = top_y + Cm(1.05) + i * row_h
        row_bg = RGBColor(0x10, 0x1E, 0x40) if i % 2 == 0 else RGBColor(0x0A, 0x14, 0x2A)
        rect(s, Cm(0.5), y, col_w * 2 + Cm(0.5), row_h - Cm(0.05), row_bg)

        # Aspect label
        txt(s, aspect, Cm(0.6), y + Cm(0.2), Cm(3.0), row_h - Cm(0.3),
            size=Pt(11), bold=True, color=ACCENT2)

        # Old value
        txt(s, old, Cm(3.5), y + Cm(0.2), Cm(11.5), row_h - Cm(0.3),
            size=Pt(11), color=RGBColor(255, 160, 160))

        # New value
        txt(s, new, mid_x + Cm(0.7), y + Cm(0.2), Cm(13.5), row_h - Cm(0.3),
            size=Pt(11), color=RGBColor(160, 255, 200))

    # Divider line
    line(s, mid_x + Cm(0.5), top_y, mid_x + Cm(0.5), H - Cm(1.2),
         ACCENT, Pt(1.5))

    footer_bar(s)


def slide_how_it_works(prs):
    """High-level flow diagram: Alice sends a message to Bob via MegaNet."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "How It Works", "Sending a private message on MegaNet — end to end")

    import math

    # Timeline/flow across the slide
    y_flow = H * 0.52
    steps = [
        (Cm(2.8),  ACCENT,  "Alice\n(Sender)",    "1. Encrypt\nmessage"),
        (Cm(8.8),  ACCENT2, "Fragment\n& Send",   "2. Split into\n203-byte chunks"),
        (Cm(15.0), GREEN,   "LoRa\nGateway",      "3. Radio waves\nover the air"),
        (Cm(21.2), ORANGE,  "Network\nNodes",     "4. Route via\nKademlia DHT"),
        (Cm(27.5), ACCENT,  "Bob\n(Receiver)",    "5. Reassemble\n& Decrypt"),
    ]

    # Draw connecting arrows between steps
    arrow_xs = [(steps[i][0] + Cm(1.0), steps[i+1][0] - Cm(1.0)) for i in range(len(steps)-1)]
    for x1, x2 in arrow_xs:
        line(s, x1, y_flow, x2, y_flow, ACCENT2, Pt(2.5))
        # Arrowhead triangle (approximate with small rect)
        rect(s, x2 - Cm(0.2), y_flow - Cm(0.15), Cm(0.2), Cm(0.3), ACCENT2)

    # Draw step circles & labels
    for cx, col, top_label, bottom_label in steps:
        circle(s, cx, y_flow, Cm(1.1), col)
        # Icon label inside circle (number)
        step_num = str(steps.index((cx, col, top_label, bottom_label)) + 1)
        txt(s, step_num, cx - Cm(0.3), y_flow - Cm(0.4), Cm(0.6), Cm(0.8),
            size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Label above
        txt(s, top_label, cx - Cm(1.5), y_flow - Cm(3.0), Cm(3.0), Cm(1.5),
            size=Pt(12), bold=True, color=col, align=PP_ALIGN.CENTER)
        # Description below
        txt(s, bottom_label, cx - Cm(1.8), y_flow + Cm(1.3), Cm(3.6), Cm(1.6),
            size=Pt(11), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # Blockchain banner at bottom
    rect(s, Cm(1), H - Cm(3.2), W - Cm(2), Cm(1.5), RGBColor(0x0F, 0x2E, 0x6A))
    txt(s, "Blockchain: every registration, routing update, and delivery receipt is anchored on-chain — immutable, verifiable, decentralised",
        Cm(1.5), H - Cm(3.1), W - Cm(3), Cm(1.3),
        size=Pt(12), color=ACCENT2, align=PP_ALIGN.CENTER, wrap=True)

    footer_bar(s)


def slide_network_map(prs):
    """Visual network topology: nodes, gateways, coverage circles."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "Network Topology", "Decentralised mesh — no single point of control or failure")

    import math

    # Canvas area (below header)
    MAP_L = Cm(1)
    MAP_T = Cm(2.7)
    MAP_W = W - Cm(2)
    MAP_H = H - Cm(3.8)

    # Background grid area
    rect(s, MAP_L, MAP_T, MAP_W, MAP_H, RGBColor(0x08, 0x12, 0x28))

    # Grid lines (subtle)
    for i in range(1, 5):
        gx = MAP_L + MAP_W * i / 5
        gy = MAP_T + MAP_H * i / 5
        line(s, gx, MAP_T, gx, MAP_T + MAP_H, RGBColor(0x1A, 0x28, 0x4A), Pt(0.5))
        line(s, MAP_L, gy, MAP_L + MAP_W, gy, RGBColor(0x1A, 0x28, 0x4A), Pt(0.5))

    # Gateway positions (5 gateways: 4 corners + centre)
    def map_pos(x_frac, y_frac):
        return (MAP_L + MAP_W * x_frac, MAP_T + MAP_H * y_frac)

    gw_positions = [
        map_pos(0.08, 0.10),   # TL
        map_pos(0.92, 0.10),   # TR
        map_pos(0.08, 0.88),   # BL
        map_pos(0.92, 0.88),   # BR
        map_pos(0.50, 0.50),   # Centre
    ]

    # Coverage circles (large, semi-transparent via layering)
    cov_radius = MAP_W * 0.28
    for gx, gy in gw_positions:
        # Outer glow rings
        for r_mult, col in [
            (1.0, RGBColor(0x00, 0x40, 0x80)),
            (0.7, RGBColor(0x00, 0x60, 0xB0)),
            (0.4, RGBColor(0x00, 0x90, 0xD0)),
        ]:
            r = cov_radius * r_mult
            shp = slide.shapes.add_shape(9,
                                          int(gx - r), int(gy - r),
                                          int(2 * r), int(2 * r)) if False else \
                  s.shapes.add_shape(9, int(gx - r), int(gy - r), int(2 * r), int(2 * r))
            shp.fill.solid()
            shp.fill.fore_color.rgb = col
            shp.line.fill.background()
            # Make it very transparent by using a very low-opacity color
            # (fpdf2 doesn't support true alpha; layer trick)

    # End-user nodes (phones/devices) — random-ish positions
    node_positions = [
        map_pos(0.20, 0.28), map_pos(0.35, 0.18), map_pos(0.15, 0.55),
        map_pos(0.28, 0.72), map_pos(0.45, 0.35), map_pos(0.52, 0.68),
        map_pos(0.65, 0.22), map_pos(0.72, 0.45), map_pos(0.80, 0.30),
        map_pos(0.60, 0.75), map_pos(0.78, 0.72), map_pos(0.40, 0.85),
        map_pos(0.55, 0.88),
    ]

    node_names = ["Alice", "Bob", "Carol", "Dave", "Eve",
                  "Frank", "Grace", "Hal", "Iris",
                  "Jules", "Karl", "Lena", "Mo"]

    # Draw mesh connections between nearby nodes
    for i, (nx, ny) in enumerate(node_positions):
        for j, (mx, my) in enumerate(node_positions):
            if j <= i:
                continue
            dist = math.sqrt((nx - mx) ** 2 + (ny - my) ** 2)
            if dist < MAP_W * 0.25:
                line(s, nx, ny, mx, my, RGBColor(0x00, 0x50, 0x90), Pt(0.8))

    # Draw gateway connections to nearest nodes
    for gx, gy in gw_positions:
        for nx, ny in node_positions:
            dist = math.sqrt((nx - gx) ** 2 + (ny - gy) ** 2)
            if dist < MAP_W * 0.35:
                line(s, gx, gy, nx, ny, RGBColor(0x00, 0x80, 0xC0), Pt(1.0))

    # Draw gateway symbols (larger circles with antenna look)
    for i, (gx, gy) in enumerate(gw_positions):
        circle(s, gx, gy, Cm(0.7), ACCENT)
        circle(s, gx, gy, Cm(0.35), WHITE)
        txt(s, "GW", gx - Cm(0.4), gy + Cm(0.75), Cm(0.8), Cm(0.5),
            size=Pt(9), bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

    # Draw user nodes
    for (nx, ny), name in zip(node_positions, node_names):
        circle(s, nx, ny, Cm(0.28), GREEN)
        txt(s, name, nx - Cm(0.7), ny + Cm(0.32), Cm(1.4), Cm(0.45),
            size=Pt(8), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # Legend
    rect(s, MAP_L + MAP_W - Cm(6.5), MAP_T + Cm(0.3), Cm(6.2), Cm(2.5),
         RGBColor(0x10, 0x20, 0x48))
    txt(s, "Legend", MAP_L + MAP_W - Cm(6.3), MAP_T + Cm(0.4), Cm(6), Cm(0.6),
        size=Pt(11), bold=True, color=WHITE)
    for dy, col, label in [
        (Cm(0.95), ACCENT,  "  LoRa Gateway (GW)"),
        (Cm(1.55), GREEN,   "  User Node (device)"),
    ]:
        circle(s, MAP_L + MAP_W - Cm(5.8), MAP_T + Cm(0.3) + dy + Cm(0.15),
               Cm(0.2), col)
        txt(s, label,
            MAP_L + MAP_W - Cm(5.5), MAP_T + Cm(0.3) + dy,
            Cm(5), Cm(0.5), size=Pt(10), color=LIGHT_GREY)

    footer_bar(s)


def slide_crypto_visual(prs):
    """Visual explanation of public-key addressing and encryption layers."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "Security Architecture", "Every bit is encrypted. Every node is self-sovereign.")

    # Left panel: Address derivation
    rect(s, Cm(0.5), Cm(2.7), Cm(15.5), Cm(7.5), RGBColor(0x0A, 0x1E, 0x44))
    txt(s, "Public-Key Node Identity", Cm(0.8), Cm(2.85), Cm(15), Cm(0.7),
        size=Pt(14), bold=True, color=ACCENT2)

    # Key chain diagram
    steps_addr = [
        (Cm(1.5), Cm(4.2), Cm(4.5), Cm(1.0), RGBColor(0x20, 0x50, 0xA0), "Ed25519\nPrivate Key"),
        (Cm(7.0), Cm(4.2), Cm(4.5), Cm(1.0), RGBColor(0x10, 0x70, 0xB0), "Ed25519\nPublic Key"),
        (Cm(12.2), Cm(4.2), Cm(3.2), Cm(1.0), ACCENT,                      "Node\nAddress"),
    ]
    for bx, by, bw, bh, bc, bl in steps_addr:
        rounded_rect(s, bx, by, bw, bh, bc)
        txt(s, bl, bx, by, bw, bh, size=Pt(11), bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

    # Arrows between
    line(s, Cm(6.0), Cm(4.7), Cm(7.0), Cm(4.7), ACCENT2, Pt(2))
    line(s, Cm(11.5), Cm(4.7), Cm(12.2), Cm(4.7), ACCENT2, Pt(2))
    txt(s, "generate", Cm(5.8), Cm(4.1), Cm(1.4), Cm(0.5), size=Pt(8), color=LIGHT_GREY)
    txt(s, "SHA3-256[:20]", Cm(11.0), Cm(4.1), Cm(1.5), Cm(0.5), size=Pt(8), color=LIGHT_GREY)

    # Address formula
    rect(s, Cm(0.8), Cm(5.5), Cm(15), Cm(0.7), RGBColor(0x05, 0x10, 0x2E))
    txt(s, "Address = SHA3-256(public_key_bytes)[:20]  ->  unique 40-hex-char ID, no registry needed",
        Cm(0.9), Cm(5.55), Cm(14.8), Cm(0.65),
        size=Pt(10), color=ACCENT2, italic=True, align=PP_ALIGN.CENTER)

    # Key properties
    props = ["No ISP registration", "No government ID", "Self-sovereign", "Cryptographically unique"]
    for i, p in enumerate(props):
        px = Cm(0.9) + i * Cm(3.8)
        circle(s, px + Cm(0.2), Cm(7.2), Cm(0.18), GREEN)
        txt(s, p, px + Cm(0.5), Cm(7.0), Cm(3.4), Cm(0.55),
            size=Pt(10), color=WHITE)

    # Right panel: Encryption layers
    rect(s, Cm(16.5), Cm(2.7), Cm(16.8), Cm(7.5), RGBColor(0x0A, 0x1E, 0x44))
    txt(s, "End-to-End Encryption Flow", Cm(16.8), Cm(2.85), Cm(16), Cm(0.7),
        size=Pt(14), bold=True, color=ACCENT2)

    enc_steps = [
        (RGBColor(0x40, 0x10, 0x10), "Plaintext message"),
        (RGBColor(0x40, 0x30, 0x10), "X25519 ECDH -> shared secret"),
        (RGBColor(0x10, 0x40, 0x10), "ChaCha20-Poly1305 encrypt"),
        (RGBColor(0x10, 0x30, 0x60), "Fragment (203B chunks)"),
        (RGBColor(0x10, 0x50, 0x80), "LoRa radio transmission"),
        (ACCENT,                       "On-chain receipt proof"),
    ]
    for i, (ec, el) in enumerate(enc_steps):
        ey = Cm(3.7) + i * Cm(0.95)
        pw = Cm(15.5) - i * Cm(0.6)
        px = Cm(16.8) + i * Cm(0.3)
        rounded_rect(s, px, ey, pw, Cm(0.85), ec)
        txt(s, el, px + Cm(0.3), ey, pw - Cm(0.4), Cm(0.85),
            size=Pt(11), color=WHITE)

    # Bottom row: threat cards
    threats = [
        (ACCENT,  "Eavesdropping", "Impossible: data\nencrypted in transit"),
        (GREEN,   "Impersonation", "Impossible: Ed25519\nsignatures required"),
        (ORANGE,  "Censorship",    "Impractical: no IP\nor DNS to block"),
        (RED,     "Tampering",     "Detectable: Poly1305\nauthentication tag"),
    ]
    for i, (tc, th, td) in enumerate(threats):
        tx = Cm(0.5) + i * Cm(8.2)
        rect(s, tx, Cm(10.6), Cm(7.8), Cm(2.5), RGBColor(0x0C, 0x1C, 0x40))
        circle(s, tx + Cm(0.6), Cm(10.6) + Cm(0.7), Cm(0.35), tc)
        txt(s, th, tx + Cm(1.2), Cm(10.65), Cm(6.4), Cm(0.65),
            size=Pt(13), bold=True, color=tc)
        txt(s, td, tx + Cm(0.3), Cm(11.3), Cm(7.2), Cm(0.9),
            size=Pt(11), color=LIGHT_GREY)

    footer_bar(s)


def slide_blockchain_visual(prs):
    """Visual blockchain: chain of blocks with icons per tx type."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "The MegaNet Blockchain", "Immutable record of identity, routing, and delivery")

    # Chain of blocks
    block_w = Cm(5.8)
    block_h = Cm(5.0)
    gap = Cm(1.5)
    start_x = Cm(0.8)
    by = Cm(3.2)

    block_data = [
        ("Block 0\nGenesis", [("Deterministic", LIGHT_GREY), ("PoW hash: 00...", ACCENT2)], NAVY),
        ("Block 1\nRegistration", [("NODE_REGISTER x5", GREEN), ("5 nodes join", LIGHT_GREY)], RGBColor(0x10, 0x30, 0x60)),
        ("Block 2\nReceipts", [("MESSAGE_RECEIPT", ORANGE), ("DATA_ANCHOR", ACCENT2)], RGBColor(0x10, 0x30, 0x60)),
        ("Block N\n...", [("Continuous", LIGHT_GREY), ("growth", LIGHT_GREY)], RGBColor(0x10, 0x28, 0x50)),
    ]

    for i, (title, items, bc) in enumerate(block_data):
        bx = start_x + i * (block_w + gap)

        # Chain link connector
        if i > 0:
            line(s, bx - gap, by + block_h / 2, bx, by + block_h / 2,
                 ACCENT, Pt(3))
            # Hash arrow label
            txt(s, "prev_hash", bx - gap - Cm(0.2), by + block_h / 2 + Cm(0.1),
                gap, Cm(0.5), size=Pt(8), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

        # Block body
        rect(s, bx, by, block_w, block_h, bc)
        # Block header bar
        rect(s, bx, by, block_w, Cm(0.9), ACCENT)
        txt(s, title, bx, by, block_w, Cm(0.9),
            size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # PoW badge
        rect(s, bx, by + Cm(0.9), block_w, Cm(0.45), RGBColor(0x05, 0x20, 0x50))
        txt(s, "SHA3-256 PoW  |  difficulty=2  |  ~1ms", bx, by + Cm(0.9), block_w, Cm(0.45),
            size=Pt(8), color=ACCENT2, align=PP_ALIGN.CENTER)

        # Transaction list
        for j, (item_text, item_color) in enumerate(items):
            iy = by + Cm(1.5) + j * Cm(0.75)
            circle(s, bx + Cm(0.4), iy + Cm(0.2), Cm(0.15), item_color)
            txt(s, item_text, bx + Cm(0.7), iy, block_w - Cm(0.8), Cm(0.65),
                size=Pt(10), color=item_color)

        # Block index badge
        circle(s, bx + block_w - Cm(0.55), by + block_h - Cm(0.55), Cm(0.45), ACCENT)
        txt(s, str(i), bx + block_w - Cm(0.9), by + block_h - Cm(0.9),
            Cm(0.7), Cm(0.7), size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # State index visual (right side)
    si_x = start_x + 4 * (block_w + gap) - Cm(0.5)
    si_y = by
    si_w = W - si_x - Cm(1)
    rect(s, si_x, si_y, si_w, block_h, RGBColor(0x08, 0x1C, 0x3C))
    txt(s, "Live State\nIndices", si_x + Cm(0.3), si_y + Cm(0.2), si_w - Cm(0.5), Cm(1.0),
        size=Pt(12), bold=True, color=ACCENT2)
    for j, (idx, col) in enumerate([
        ("node_registry", GREEN),
        ("message_receipts", ORANGE),
        ("data_anchors", ACCENT),
    ]):
        idy = si_y + Cm(1.4) + j * Cm(1.1)
        rounded_rect(s, si_x + Cm(0.3), idy, si_w - Cm(0.6), Cm(0.9), col)
        txt(s, idx, si_x + Cm(0.5), idy, si_w - Cm(0.8), Cm(0.9),
            size=Pt(11), bold=True, color=WHITE)

    # Bottom: properties
    props = [
        (ACCENT2, "Immutable", "Once written, blocks cannot be altered without redoing all PoW"),
        (GREEN,   "Verifiable", "Every transaction is signed; anyone can verify"),
        (ORANGE,  "Decentralised", "Every node holds a full copy — no central server"),
    ]
    for i, (pc, pt, pd) in enumerate(props):
        px = Cm(0.5) + i * Cm(11.1)
        rect(s, px, H - Cm(4.0), Cm(10.7), Cm(2.5), RGBColor(0x0C, 0x1C, 0x40))
        txt(s, pt, px + Cm(0.4), H - Cm(3.9), Cm(10), Cm(0.7),
            size=Pt(14), bold=True, color=pc)
        txt(s, pd, px + Cm(0.4), H - Cm(3.1), Cm(10), Cm(1.5),
            size=Pt(11), color=LIGHT_GREY, wrap=True)

    footer_bar(s)


def slide_lora_coverage(prs):
    """LoRa range & SF chart + coverage diagram."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "LoRaWAN Radio Transport", "Long-range, low-power — no ISP, no cables needed")

    # Left: Coverage diagram
    MAP_L = Cm(0.5)
    MAP_T = Cm(2.7)
    MAP_W = Cm(14)
    MAP_H = H - Cm(3.8)

    rect(s, MAP_L, MAP_T, MAP_W, MAP_H, RGBColor(0x08, 0x12, 0x28))
    txt(s, "LoRa Coverage (20x20 km grid)", MAP_L, MAP_T + Cm(0.1), MAP_W, Cm(0.5),
        size=Pt(11), bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

    # Gateway in centre
    cx = MAP_L + MAP_W / 2
    cy = MAP_T + MAP_H / 2

    # Concentric coverage rings (SF7 to SF12)
    sf_rings = [
        (0.32, RGBColor(0x80, 0x20, 0x00), "SF7 ~7km"),
        (0.44, RGBColor(0x80, 0x50, 0x00), "SF9 ~11km"),
        (0.56, RGBColor(0x00, 0x60, 0x20), "SF12 ~17km"),
    ]
    max_r = min(MAP_W, MAP_H) * 0.52
    for frac, col, label in sf_rings:
        r = max_r * frac
        shp = s.shapes.add_shape(9, int(cx - r), int(cy - r), int(2 * r), int(2 * r))
        shp.fill.solid()
        shp.fill.fore_color.rgb = col
        r_val, g_val, b_val = col[0], col[1], col[2]
        shp.line.color.rgb = RGBColor(min(r_val + 60, 255),
                                       min(g_val + 60, 255),
                                       min(b_val + 60, 255))
        shp.line.width = Pt(1)
        # Label
        txt(s, label, cx + r * 0.65, cy - Cm(0.3), Cm(3), Cm(0.5),
            size=Pt(9), color=WHITE)

    # Gateway tower
    circle(s, cx, cy, Cm(0.6), ACCENT)
    txt(s, "GW", cx - Cm(0.4), cy + Cm(0.65), Cm(0.8), Cm(0.5),
        size=Pt(10), bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

    # A few device nodes
    import math
    dev_positions = [(0.28, 0.30), (0.60, 0.22), (0.20, 0.65),
                     (0.70, 0.70), (0.80, 0.45)]
    for fx, fy in dev_positions:
        dx = MAP_L + MAP_W * fx
        dy = MAP_T + MAP_H * fy
        line(s, cx, cy, dx, dy, RGBColor(0x00, 0x70, 0xA0), Pt(1.0))
        circle(s, dx, dy, Cm(0.25), GREEN)

    # Scale bar
    bar_y = MAP_T + MAP_H - Cm(0.8)
    bar_w = MAP_W * 0.3
    line(s, MAP_L + Cm(0.5), bar_y, MAP_L + Cm(0.5) + bar_w, bar_y, WHITE, Pt(1.5))
    txt(s, "~10 km", MAP_L + Cm(0.5) + bar_w / 2 - Cm(1), bar_y - Cm(0.55),
        Cm(2), Cm(0.5), size=Pt(9), color=WHITE, align=PP_ALIGN.CENTER)

    # Right: SF parameter chart
    rect(s, Cm(15), MAP_T, W - Cm(15.5), MAP_H, RGBColor(0x08, 0x12, 0x28))
    txt(s, "Spreading Factor Trade-offs", Cm(15.2), MAP_T + Cm(0.15), Cm(17), Cm(0.5),
        size=Pt(11), bold=True, color=ACCENT2)

    # Bar chart: Range vs Data Rate
    chart_data = ChartData()
    chart_data.categories = ["SF7", "SF8", "SF9", "SF10", "SF11", "SF12"]
    chart_data.add_series("Range (km)", [7, 9, 11, 13, 15, 17])
    chart_data.add_series("Data rate x10 (bps)", [55, 32, 18, 11, 6, 3])

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Cm(15), MAP_T + Cm(0.8),
        W - Cm(15.8), MAP_H - Cm(0.9),
        chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = True

    plot = chart.plots[0]
    colors_series = [ACCENT, GREEN]
    for i, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors_series[i]

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = WHITE
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.color.rgb = WHITE
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = WHITE

    footer_bar(s)


def slide_use_cases(prs):
    """Executive use cases — business value."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "Use Cases & Market Opportunity", "Where MegaNet delivers unique value")

    cases = [
        (Cm(0.5),  Cm(2.8),  ACCENT,  "Disaster\nResponse",
         "When traditional networks fail,\nMegaNet keeps communicating.\nNo towers, no cables needed.",
         "Emergency Services\nNGOs  |  Military"),
        (Cm(8.9),  Cm(2.8),  GREEN,   "Censored\nRegions",
         "Citizens in authoritarian states\ncommunicate freely — no IP\nto block, no DNS to poison.",
         "Journalists\nActivists  |  Dissidents"),
        (Cm(17.3), Cm(2.8),  ORANGE,  "IoT &\nSmart Cities",
         "Billions of sensors reporting\nwithout SIM cards or ISPs.\nLoRa covers 10+ km per gateway.",
         "Agriculture\nInfrastructure  |  Logistics"),
        (Cm(25.7), Cm(2.8),  YELLOW,  "Privacy-\nFirst Apps",
         "Messaging, payments, and\nidentity without revealing\nIP address or personal data.",
         "Financial Privacy\nSecure Comms  |  Web3"),

        (Cm(0.5),  Cm(9.8),  RGBColor(0x80,0x40,0xC0), "Offshore &\nMaritime",
         "Ships and remote platforms\nwith LoRa coverage over\nhundreds of km of open sea.",
         "Shipping\nFishing  |  Energy"),
        (Cm(8.9),  Cm(9.8),  ACCENT2, "National\nSecurity",
         "Resilient backup comms\nthat survive infrastructure\nattacks or EMP events.",
         "Government\nDefence  |  Intelligence"),
        (Cm(17.3), Cm(9.8),  GREEN,   "Rural\nConnectivity",
         "Bring internet-like services\nto areas with no mobile\ncoverage at near-zero cost.",
         "Rural Communities\nDeveloping Nations"),
        (Cm(25.7), Cm(9.8),  RED,     "Critical\nInfrastructure",
         "Power grids, water systems,\ntransport networks with\ncensorship-proof control planes.",
         "Utilities\nTransport  |  Industry"),
    ]

    for cx, cy, col, title, body, market in cases:
        card_w = Cm(7.8)
        card_h = Cm(6.2)
        rect(s, cx, cy, card_w, card_h, RGBColor(0x0C, 0x1C, 0x40))
        rect(s, cx, cy, card_w, Cm(0.5), col)
        circle(s, cx + card_w - Cm(0.7), cy + Cm(0.8), Cm(0.45), col)
        txt(s, title, cx + Cm(0.3), cy + Cm(0.55), Cm(5.5), Cm(1.2),
            size=Pt(14), bold=True, color=col)
        txt(s, body, cx + Cm(0.3), cy + Cm(1.8), card_w - Cm(0.6), Cm(2.5),
            size=Pt(10.5), color=WHITE, wrap=True)
        txt(s, market, cx + Cm(0.3), cy + Cm(4.5), card_w - Cm(0.6), Cm(0.9),
            size=Pt(9.5), color=col, italic=True, wrap=True)

    footer_bar(s)


def slide_metrics_dashboard(prs):
    """Dashboard slide: key performance metrics with visual KPI cards."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "Key Metrics & Performance", "What the MVP achieves — measured, verified")

    # Top row: 4 big KPI numbers
    kpis = [
        (Cm(0.5),  "42",       "Unit Tests\nPassing",        ACCENT,  "100% pass rate"),
        (Cm(8.7),  "0.16s",    "Full Test\nSuite Runtime",   GREEN,   "No hardware needed"),
        (Cm(16.9), "~1ms",     "Block\nMine Time",           ORANGE,  "PoW difficulty=2"),
        (Cm(25.1), "203 B",    "Max Fragment\nSize (SF7)",   ACCENT2, "19B header overhead"),
    ]

    for kx, num, label, col, sub in kpis:
        kw = Cm(7.8)
        kh = Cm(4.0)
        rect(s, kx, Cm(2.7), kw, kh, RGBColor(0x0C, 0x1C, 0x40))
        rect(s, kx, Cm(2.7), kw, Cm(0.18), col)
        txt(s, num, kx, Cm(3.1), kw, Cm(1.8),
            size=Pt(44), bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, label, kx, Cm(4.9), kw, Cm(0.9),
            size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, sub, kx, Cm(5.8), kw, Cm(0.6),
            size=Pt(10), color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

    # Bottom left: bar chart — RSSI vs distance
    chart_data = ChartData()
    chart_data.categories = ["1km", "5km", "10km", "15km"]
    chart_data.add_series("RSSI (dBm + 150 offset)", [52, 33, 25, 20])

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Cm(0.5), Cm(7.0),
        Cm(14), Cm(9.0),
        chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Signal Strength vs Distance (EU868)"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = WHITE
    chart.has_legend = False
    plot = chart.plots[0]
    for series in plot.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = ACCENT
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.category_axis.tick_labels.font.color.rgb = WHITE
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.color.rgb = WHITE

    # Annotation
    txt(s, "SF7 sensitivity: -123 dBm\n(-98 dBm @ 1km = +25 dBm margin!)",
        Cm(0.8), Cm(15.1), Cm(12), Cm(0.9),
        size=Pt(10), color=ACCENT2, italic=True)

    # Bottom right: delivery probability chart
    chart_data2 = ChartData()
    chart_data2.categories = ["SF7", "SF8", "SF9", "SF10", "SF11", "SF12"]
    chart_data2.add_series("Delivery % @ 5km", [85, 91, 96, 99, 99.5, 99.8])

    chart2 = s.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Cm(15.5), Cm(7.0),
        Cm(17.5), Cm(9.0),
        chart_data2
    ).chart
    chart2.has_title = True
    chart2.chart_title.text_frame.text = "Delivery Probability @ 5km by Spreading Factor"
    chart2.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart2.chart_title.text_frame.paragraphs[0].font.color.rgb = WHITE
    chart2.has_legend = False
    plot2 = chart2.plots[0]
    for series in plot2.series:
        series.format.line.color.rgb = GREEN
        series.format.line.width = Pt(2.5)
    chart2.category_axis.tick_labels.font.size = Pt(11)
    chart2.category_axis.tick_labels.font.color.rgb = WHITE
    chart2.value_axis.tick_labels.font.size = Pt(10)
    chart2.value_axis.tick_labels.font.color.rgb = WHITE

    footer_bar(s)


def slide_roadmap(prs):
    """Visual roadmap timeline."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "Roadmap", "From MVP prototype to production global network")

    # Timeline line
    tl_y = H * 0.48
    tl_start = Cm(1.5)
    tl_end = W - Cm(1.5)
    line(s, tl_start, tl_y, tl_end, tl_y, ACCENT, Pt(3))

    phases = [
        (0.06,  "Q1 2026",   "MVP",              "42 tests\nPure Python\nSimulation",
         ACCENT, "DONE", True),
        (0.28,  "Q3 2026",   "Hardware",         "SX1276 driver\nRaspberry Pi\nESP32 nodes",
         GREEN,  "IN PROGRESS", False),
        (0.50,  "Q1 2027",   "Protocol\nHarden", "Key revocation\nPoW scale\nStore-forward",
         ORANGE, "PLANNED", False),
        (0.72,  "Q3 2027",   "Incentive\nLayer", "Token rewards\nProof-of-BW\nEconomy",
         YELLOW, "PLANNED", False),
        (0.94,  "2028+",     "Global\nNetwork",  "App layer\nBrowser ext\nMobile app",
         ACCENT2,"VISION",  False),
    ]

    for frac, date, phase_title, details, col, status, done in phases:
        px = tl_start + (tl_end - tl_start) * frac

        # Dot on timeline
        circle(s, px, tl_y, Cm(0.5), col)
        if done:
            txt(s, "v", px - Cm(0.2), tl_y - Cm(0.35), Cm(0.4), Cm(0.5),
                size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Alternating above/below
        above = (phases.index((frac, date, phase_title, details, col, status, done)) % 2 == 0)
        box_h = Cm(3.8)
        if above:
            box_y = tl_y - Cm(0.6) - box_h
            line(s, px, box_y + box_h, px, tl_y - Cm(0.5), col, Pt(1.5))
        else:
            box_y = tl_y + Cm(0.6)
            line(s, px, tl_y + Cm(0.5), px, box_y, col, Pt(1.5))

        box_w = Cm(5.8)
        rect(s, px - box_w / 2, box_y, box_w, box_h, RGBColor(0x0C, 0x1C, 0x40))
        rect(s, px - box_w / 2, box_y, box_w, Cm(0.45), col)

        txt(s, date, px - box_w / 2, box_y, box_w, Cm(0.45),
            size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, phase_title, px - box_w / 2, box_y + Cm(0.5), box_w, Cm(0.9),
            size=Pt(14), bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, details, px - box_w / 2, box_y + Cm(1.45), box_w, Cm(1.6),
            size=Pt(10), color=LIGHT_GREY, align=PP_ALIGN.CENTER, wrap=True)

        # Status badge
        rect(s, px - box_w / 2 + Cm(0.4), box_y + box_h - Cm(0.7),
             box_w - Cm(0.8), Cm(0.55), col)
        txt(s, status, px - box_w / 2 + Cm(0.4), box_y + box_h - Cm(0.7),
            box_w - Cm(0.8), Cm(0.55),
            size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer_bar(s)


def slide_competitive(prs):
    """Competitive landscape radar / comparison."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "Competitive Landscape", "How MegaNet compares to existing solutions")

    # Comparison matrix
    competitors = ["Traditional\nInternet", "Tor / VPN", "Mesh\nNetworks", "Satellite\n(Starlink)", "MegaNet"]
    attributes = [
        ("Censorship Resistance", [1, 3, 3, 2, 5]),
        ("Privacy by Default",   [1, 4, 2, 2, 5]),
        ("No Infrastructure",    [1, 1, 3, 2, 5]),
        ("Low Cost",             [3, 4, 4, 2, 5]),
        ("Long Range",           [5, 5, 2, 5, 4]),
        ("Data Rate",            [5, 5, 3, 4, 2]),
        ("Decentralised",        [1, 2, 3, 2, 5]),
        ("Open Protocol",        [2, 3, 4, 2, 5]),
    ]

    # Table headers
    tbl_l = Cm(0.5)
    tbl_t = Cm(2.7)
    col_w_first = Cm(7.5)
    col_w = Cm(4.8)
    row_h = Cm(1.25)

    # Header row
    rect(s, tbl_l, tbl_t, col_w_first, row_h, NAVY)
    for i, comp in enumerate(competitors):
        cx = tbl_l + col_w_first + i * col_w
        hcol = ACCENT if comp == "MegaNet" else NAVY
        rect(s, cx, tbl_t, col_w - Cm(0.05), row_h, hcol)
        txt(s, comp, cx, tbl_t, col_w - Cm(0.05), row_h,
            size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Data rows
    for r, (attr, scores) in enumerate(attributes):
        ry = tbl_t + (r + 1) * row_h
        bg_row = RGBColor(0x10, 0x20, 0x44) if r % 2 == 0 else RGBColor(0x0A, 0x14, 0x34)
        rect(s, tbl_l, ry, col_w_first + len(competitors) * col_w, row_h - Cm(0.05), bg_row)

        txt(s, attr, tbl_l + Cm(0.2), ry, col_w_first - Cm(0.3), row_h - Cm(0.05),
            size=Pt(11), color=WHITE)

        for i, (score, comp) in enumerate(zip(scores, competitors)):
            cx = tbl_l + col_w_first + i * col_w
            # Score dots
            is_mega = (comp == "MegaNet")
            for dot_i in range(5):
                filled = dot_i < score
                dc = GREEN if (filled and is_mega) else (ACCENT if filled else MID_GREY)
                circle(s, cx + Cm(0.5) + dot_i * Cm(0.75), ry + row_h / 2,
                       Cm(0.22) if filled else Cm(0.18), dc)

    # Legend
    rect(s, tbl_l, tbl_t + 9 * row_h + Cm(0.3), Cm(12), Cm(0.7), RGBColor(0x0A, 0x14, 0x34))
    circle(s, tbl_l + Cm(0.4), tbl_t + 9 * row_h + Cm(0.65), Cm(0.2), GREEN)
    txt(s, "Filled = capability present  |  5 dots = maximum  |  Green = MegaNet advantage",
        tbl_l + Cm(0.8), tbl_t + 9 * row_h + Cm(0.35), Cm(11), Cm(0.6),
        size=Pt(10), color=LIGHT_GREY)

    footer_bar(s)


def slide_summary(prs):
    """Executive summary — 3 key takeaways."""
    s = new_slide(prs)
    bg(s)
    header_bar(s, "Executive Summary", "Three reasons MegaNet matters")

    pillars = [
        (Cm(0.5),  ACCENT,  "1. Freedom",
         "No IP addresses means no government or corporation\n"
         "can switch off your ability to communicate.\n\n"
         "MegaNet is the first protocol where the address\n"
         "is mathematically derived from who you are,\n"
         "not where you live or who your ISP is."),
        (Cm(11.5), GREEN,   "2. Security",
         "Every message is encrypted end-to-end using\n"
         "best-in-class algorithms (X25519 + ChaCha20).\n\n"
         "Every transaction is signed and immutably\n"
         "recorded on a blockchain — creating a permanent,\n"
         "tamper-proof audit trail."),
        (Cm(22.5), ORANGE,  "3. Resilience",
         "No single gateway, server, or nation can take\n"
         "down the network. Data is fragmented and\n"
         "distributed — 'homeless', unownable.\n\n"
         "LoRa radio links survive when fibre is cut,\n"
         "towers are destroyed, or ISPs are coerced."),
    ]

    for px, col, title, body in pillars:
        pw = Cm(10.5)
        ph = H - Cm(6.2)
        rect(s, px, Cm(2.7), pw, ph, RGBColor(0x0C, 0x1C, 0x40))
        rect(s, px, Cm(2.7), pw, Cm(0.5), col)
        circle(s, px + pw / 2, Cm(2.7) + Cm(1.5), Cm(0.9), col)
        txt(s, title.split(". ")[0], px + pw / 2 - Cm(0.5), Cm(2.7) + Cm(0.9),
            Cm(1), Cm(1.0), size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, px + Cm(0.4), Cm(2.7) + Cm(2.8), pw - Cm(0.8), Cm(1.0),
            size=Pt(17), bold=True, color=col)
        txt(s, body, px + Cm(0.4), Cm(2.7) + Cm(3.9), pw - Cm(0.8), ph - Cm(4.2),
            size=Pt(12), color=WHITE, wrap=True)

    # Bottom call to action
    rect(s, Cm(0), H - Cm(2.5), W, Cm(1.6), ACCENT)
    txt(s,
        "MVP proven  |  42/42 tests pass  |  Full E2E demo: python main.py  |  Ready for hardware phase",
        Cm(1), H - Cm(2.4), W - Cm(2), Cm(1.4),
        size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer_bar(s)


def slide_closing(prs):
    s = new_slide(prs)
    bg(s, DARK_BG)

    # Background accent column
    rect(s, W * 0.68, 0, W * 0.32, H, RGBColor(0x05, 0x0E, 0x22))
    for i in range(5):
        c = RGBColor(0, max(0, 60 - i * 10), min(255, 120 + i * 25))
        rect(s, W * 0.68 + Cm(i * 1.1), 0, Cm(0.4), H, c)

    # Central node network (right side visual)
    import math
    cx2, cy2 = W * 0.84, H * 0.45
    circle(s, cx2, cy2, Cm(3.5), RGBColor(0, 25, 60))
    circle(s, cx2, cy2, Cm(2.2), RGBColor(0, 50, 110))
    circle(s, cx2, cy2, Cm(1.2), ACCENT)
    circle(s, cx2, cy2, Cm(0.5), WHITE)
    for angle, nc in [(45, GREEN), (135, ORANGE), (225, ACCENT2), (315, YELLOW)]:
        nx = cx2 + Cm(4.0) * math.cos(math.radians(angle))
        ny = cy2 + Cm(4.0) * math.sin(math.radians(angle))
        line(s, cx2, cy2, nx, ny, ACCENT, Pt(1.5))
        circle(s, nx, ny, Cm(0.45), nc)

    multi_para(s, [
        ("The internet", Pt(52), WHITE, True, PP_ALIGN.LEFT),
        ("can be rebuilt.", Pt(52), ACCENT, True, PP_ALIGN.LEFT),
        ("", Pt(16), WHITE, False, PP_ALIGN.LEFT),
        ("No IPs. No DNS. No censorship.", Pt(20), ACCENT2, False, PP_ALIGN.LEFT),
        ("No single point of failure.", Pt(20), ACCENT2, False, PP_ALIGN.LEFT),
        ("", Pt(14), WHITE, False, PP_ALIGN.LEFT),
        ("MegaNet MVP v1.0 | 2026", Pt(13), LIGHT_GREY, False, PP_ALIGN.LEFT),
    ], Cm(1.5), Cm(3.5), Cm(21), Cm(11))

    rect(s, Cm(0), H - Cm(2.5), Cm(23), Cm(2.5), ACCENT)
    txt(s, "Code: /usr/src/meganet  |  Run: python main.py  |  Questions?",
        Cm(0.5), H - Cm(2.3), Cm(22), Cm(2.0),
        size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer_bar(s, "MegaNet MVP  |  Executive Briefing  |  Hacking Conference 2026")


# =============================================================================
# MAIN
# =============================================================================
def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Generating executive slides...")
    slide_cover(prs)            ; print("  1/13 Cover")
    slide_problem(prs)          ; print("  2/13 Problem")
    slide_vision(prs)           ; print("  3/13 Vision")
    slide_how_it_works(prs)     ; print("  4/13 How It Works")
    slide_network_map(prs)      ; print("  5/13 Network Map")
    slide_crypto_visual(prs)    ; print("  6/13 Crypto")
    slide_blockchain_visual(prs); print("  7/13 Blockchain")
    slide_lora_coverage(prs)    ; print("  8/13 LoRa Coverage")
    slide_use_cases(prs)        ; print("  9/13 Use Cases")
    slide_metrics_dashboard(prs); print(" 10/13 Metrics")
    slide_roadmap(prs)          ; print(" 11/13 Roadmap")
    slide_competitive(prs)      ; print(" 12/13 Competitive")
    slide_summary(prs)          ; print(" 13/13 Summary + Closing")
    slide_closing(prs)          ; print(" 14/14 Closing")

    out = "/usr/src/meganet/MegaNet_Executive.pptx"
    prs.save(out)
    print(f"\nSaved: {out}")


if __name__ == "__main__":
    build()
