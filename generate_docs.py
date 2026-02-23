#!/usr/bin/env python3
"""Generate MegaNet technical PDF report and conference PPT."""

# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
from fpdf import FPDF
from fpdf.enums import XPos, YPos


class MegaNetPDF(FPDF):
    TITLE_COLOR  = (15, 30, 60)
    ACCENT_COLOR = (0, 120, 215)
    TEXT_COLOR   = (30, 30, 30)
    LIGHT_BG     = (235, 242, 255)

    def header(self):
        self.set_fill_color(*self.TITLE_COLOR)
        self.rect(0, 0, 210, 12, "F")
        self.set_text_color(255, 255, 255)
        self.set_font("Helvetica", "B", 9)
        self.set_xy(10, 2)
        self.cell(0, 8, "MegaNet MVP - Technical Report  |  Internet Paralela Descentralizada",
                  new_x=XPos.RIGHT, new_y=YPos.TOP)
        self.ln(14)

    def footer(self):
        self.set_y(-12)
        self.set_fill_color(*self.TITLE_COLOR)
        self.rect(0, 285, 210, 12, "F")
        self.set_text_color(255, 255, 255)
        self.set_font("Helvetica", "", 8)
        self.set_x(10)
        self.cell(0, 8, f"Page {self.page_no()}  |  2026-02-23  |  MegaNet MVP v1.0", align="C")

    def section(self, title: str) -> None:
        self.ln(4)
        self.set_fill_color(*self.ACCENT_COLOR)
        self.set_text_color(255, 255, 255)
        self.set_font("Helvetica", "B", 11)
        self.cell(0, 8, f"  {title}", new_x=XPos.LMARGIN, new_y=YPos.NEXT, fill=True)
        self.set_text_color(*self.TEXT_COLOR)
        self.ln(2)

    def subsection(self, title: str) -> None:
        self.set_text_color(*self.ACCENT_COLOR)
        self.set_font("Helvetica", "B", 10)
        self.cell(0, 7, title, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        self.set_text_color(*self.TEXT_COLOR)

    def body(self, text: str) -> None:
        self.set_font("Helvetica", "", 9)
        self.set_text_color(*self.TEXT_COLOR)
        self.multi_cell(0, 5, text)
        self.ln(1)

    def code(self, text: str) -> None:
        self.set_fill_color(*self.LIGHT_BG)
        self.set_text_color(20, 20, 80)
        self.set_font("Courier", "", 7.5)
        self.multi_cell(0, 4.5, text, fill=True)
        self.set_text_color(*self.TEXT_COLOR)
        self.ln(1)

    def table(self, headers: list, rows: list, col_widths: list) -> None:
        # Header row
        self.set_fill_color(*self.ACCENT_COLOR)
        self.set_text_color(255, 255, 255)
        self.set_font("Helvetica", "B", 8)
        for h, w in zip(headers, col_widths):
            self.cell(w, 6, h, border=1, fill=True,
                      new_x=XPos.RIGHT, new_y=YPos.TOP)
        self.ln()
        # Data rows
        self.set_text_color(*self.TEXT_COLOR)
        self.set_font("Helvetica", "", 8)
        for i, row in enumerate(rows):
            if i % 2 == 0:
                self.set_fill_color(245, 248, 255)
            else:
                self.set_fill_color(255, 255, 255)
            for cell, w in zip(row, col_widths):
                self.cell(w, 5.5, str(cell), border=1, fill=True,
                          new_x=XPos.RIGHT, new_y=YPos.TOP)
            self.ln()
        self.ln(2)


def build_pdf():
    pdf = MegaNetPDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()

    # Cover
    pdf.set_fill_color(*MegaNetPDF.TITLE_COLOR)
    pdf.set_text_color(255, 255, 255)
    pdf.set_font("Helvetica", "B", 24)
    pdf.ln(6)
    pdf.cell(0, 14, "MegaNet MVP", new_x=XPos.LMARGIN, new_y=YPos.NEXT, align="C")
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 8, "Internet Paralela Descentralizada",
             new_x=XPos.LMARGIN, new_y=YPos.NEXT, align="C")
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(100, 160, 220)
    pdf.cell(0, 6, "Technical Report  |  v1.0  |  2026-02-23",
             new_x=XPos.LMARGIN, new_y=YPos.NEXT, align="C")
    pdf.set_text_color(*MegaNetPDF.TEXT_COLOR)
    pdf.ln(6)

    # 1. Executive Summary
    pdf.section("1. Executive Summary")
    pdf.body(
        "MegaNet is a proof-of-concept of Kim Dotcom's vision of a parallel, "
        "decentralised internet that does not rely on the IP address space. "
        "Nodes are identified by public keys; communication is encrypted end-to-end; "
        "all state is anchored in a lightweight blockchain; "
        "and physical transport is simulated with an EU868 LoRaWAN model.\n\n"
        "Key properties: public-key addressing (SHA3-256[:20] = 160-bit), "
        "X25519 ECDH + ChaCha20-Poly1305 encryption, SHA3-256 PoW blockchain, "
        "Kademlia DHT routing, 'homeless data' content-addressed fragments."
    )

    # 2. Architecture
    pdf.section("2. System Architecture")
    pdf.body("The system is organised into five layers, each independently testable:")
    pdf.table(
        ["Layer", "Module", "Responsibility"],
        [
            ["Crypto",       "crypto/keys.py",
             "Ed25519 identity, X25519 ECDH, ChaCha20-Poly1305"],
            ["Blockchain",   "blockchain/{tx,block,chain}.py",
             "PoW chain, 4 TX types, state indices, fork resolution"],
            ["Radio (LoRa)", "lora/{packet,gateway,simulator}.py",
             "EU868 propagation model, duty-cycle, 5 gateways"],
            ["Routing",      "routing/{table,dht}.py",
             "Kademlia k-buckets (160-bit XOR), fragmentation, ContentStore"],
            ["Node",         "node/node.py + network/simulator.py",
             "Integration, send/receive pipeline, churn, partitions"],
        ],
        [30, 62, 92],
    )

    # 3. Crypto
    pdf.section("3. Cryptographic Design")
    pdf.subsection("3.1 Node Address Derivation")
    pdf.code(
        "Ed25519.generate() -> (ed_priv, ed_pub)\n"
        "X25519.generate()  -> (x_priv,  x_pub)\n"
        "address = SHA3_256(ed_pub_raw_bytes)[:20]   # 160-bit, 40-char hex\n"
        "Published on-chain via NODE_REGISTER transaction"
    )
    pdf.subsection("3.2 Message Encryption (A -> B)")
    pdf.code(
        "shared_secret = X25519(A.x25519_priv, B.x25519_pub)   # 32 bytes\n"
        "aad           = A.addr + B.addr                         # 40 bytes\n"
        "nonce         = os.urandom(12)                          # 96-bit\n"
        "ct            = ChaCha20Poly1305.encrypt(shared_secret, plaintext, aad, nonce)\n"
        "wire          = A.addr(20B) + nonce(12B) + ct"
    )
    pdf.body(
        "ChaCha20-Poly1305 is chosen for IoT/LoRa suitability: "
        "hardware-independent constant-time implementation, "
        "256-bit security, and native AEAD authentication without separate HMAC."
    )

    # 4. Blockchain
    pdf.section("4. Lightweight Blockchain")
    pdf.subsection("4.1 Block Structure & PoW")
    pdf.code(
        "Block.hash = SHA3-256(JSON({index, timestamp, txs, prev_hash, miner, difficulty, nonce}))\n"
        "Mine: increment nonce until hash.startswith('00')  # difficulty=2 -> ~256 iters, ~1ms"
    )
    pdf.subsection("4.2 Transaction Types")
    pdf.table(
        ["Type", "Purpose", "Key Payload Fields"],
        [
            ["NODE_REGISTER",   "Publish public keys",
             "ed25519_pub, x25519_pub"],
            ["ROUTING_UPDATE",  "Announce known peers",
             "peers[] (max 10)"],
            ["MESSAGE_RECEIPT", "Non-repudiable delivery proof",
             "msg_id, receiver"],
            ["DATA_ANCHOR",     "Pin content hash on-chain",
             "content_hash, frag_count, ttl"],
        ],
        [42, 60, 82],
    )
    pdf.subsection("4.3 Mempool")
    pdf.body(
        "Max 50 transactions. FIFO eviction when full. "
        "Deduplication by tx_id = SHA3-256(canonical JSON)."
    )

    # 5. LoRa
    pdf.section("5. LoRa Radio Simulation (EU868)")
    pdf.subsection("5.1 Propagation Model")
    pdf.code(
        "PL(d) = PL0 + 10 * n * log10(d / d0)\n"
        "PL0 = 112 dB  (practical 1km reference, includes 6dB system margin)\n"
        "n   = 2.7     (rural/suburban path-loss exponent)\n"
        "RSSI(d) = 14 dBm - PL(d)\n\n"
        "P(delivery | RSSI, SF) = 1 / (1 + exp(-0.5 * (RSSI - sensitivity[SF])))"
    )
    pdf.subsection("5.2 SF Parameters")
    pdf.table(
        ["SF", "Data Bytes", "Sensitivity", "RSSI @ 5km", "Range"],
        [
            ["SF7",  "203 B", "-123 dBm", "-117 dBm", "~7 km"],
            ["SF8",  "203 B", "-126 dBm", "-117 dBm", "~9 km"],
            ["SF9",  "96 B",  "-129 dBm", "-117 dBm", "~11 km"],
            ["SF10", "32 B",  "-132 dBm", "-117 dBm", "~13 km"],
            ["SF12", "32 B",  "-137 dBm", "-117 dBm", "~17 km"],
        ],
        [20, 26, 28, 28, 32],
    )
    pdf.body(
        "5 gateways cover a 20x20 km grid (4 corners + centre). "
        "EU868 duty cycle: 1% per hour (36,000 ms budget), "
        "enforced by a sliding-window DutyCycleTracker at each gateway."
    )

    # 6. Kademlia DHT
    pdf.section("6. Kademlia DHT Routing")
    pdf.body(
        "Standard 160-bucket Kademlia with k=20. Bucket index = "
        "(our_addr XOR their_addr).bit_length() - 1. "
        "Contacts ordered by recency (LRS at front, MRS at back).\n\n"
        "Content addressing ('homeless data'):\n"
        "  content_hash = SHA3-256(data)[:20]\n"
        "  store_key    = SHA3-256(content_hash + frag_idx)[:20]\n"
        "Fragments are distributed across all reachable nodes; "
        "no single node owns the data."
    )

    # 7. Message Flow
    pdf.section("7. End-to-End Message Flow")
    pdf.code(
        "1. Lookup B.x25519_pub in blockchain.node_registry\n"
        "2. shared_secret = X25519(A.x25519_priv, B.x25519_pub)\n"
        "3. (nonce, ct) = ChaCha20-Poly1305(shared_secret, plaintext, aad=A.addr+B.addr)\n"
        "4. wire = A.addr(20B) + nonce(12B) + ct\n"
        "5. fragment_message(wire, max=203B) -> (content_hash, [Fragment...])\n"
        "6. for frag in fragments:\n"
        "     pkt = LoRaPacket(msg_id, frag_idx, total, SF7, frag.data)\n"
        "     LoRaSimulator.transmit(pkt, A.x_km, A.y_km)  # path loss + random\n"
        "7. B.receive_packets(collected) -> reassembly -> decrypt -> B.inbox\n"
        "8. B creates MESSAGE_RECEIPT tx -> mempool -> mine_block() -> broadcast"
    )

    # 8. Security
    pdf.section("8. Security Analysis")
    pdf.table(
        ["Threat", "Mitigation"],
        [
            ["Eavesdropping",     "ChaCha20-Poly1305 E2E encryption"],
            ["Impersonation",     "Ed25519 signatures on all transactions"],
            ["Replay attacks",    "Random 96-bit nonce per message"],
            ["Payload tampering", "Poly1305 16-byte authentication tag"],
            ["AAD manipulation",  "Sender+receiver bound in AEAD AAD"],
            ["IP censorship",     "No IP addresses in the protocol"],
            ["Sybil (partial)",   "PoW registration cost on blockchain"],
        ],
        [70, 114],
    )

    # 9. Test Results
    pdf.section("9. Test Results")
    pdf.table(
        ["Test Module", "Tests", "Areas Covered"],
        [
            ["test_crypto.py",     "11",
             "ECDH symmetry, encrypt/decrypt round-trip, tamper detection, Ed25519"],
            ["test_blockchain.py", "17",
             "Genesis, PoW validity, mempool eviction, fork resolution, state indices"],
            ["test_lora.py",       "14",
             "Packet serialisation, path loss, duty cycle, sigmoid prob, dedup"],
        ],
        [55, 18, 111],
    )
    pdf.body("42 / 42 tests pass in ~0.16 seconds on Python 3.13. No hardware required.")

    # 10. Parameters
    pdf.section("10. Key Technical Parameters")
    pdf.table(
        ["Parameter", "Value", "Justification"],
        [
            ["Fragment size",     "203 B",      "SF7 payload (222B) - 19B header"],
            ["Duty cycle",        "1% / 1 hr",  "EU868 ISM band regulation"],
            ["PoW difficulty",    "2 hex zeros", "~256 iters, ~1ms (demo speed)"],
            ["Block size limit",  "4096 B",     "~20 SF7 fragments per broadcast"],
            ["Address space",     "160 bits",   "Kademlia standard"],
            ["Path-loss exp.",    "2.7",        "Rural/suburban EU868"],
            ["PL ref @ 1km",      "112 dB",     "Practical LoRa link budget"],
            ["k-bucket size",     "20",         "Kademlia standard k"],
        ],
        [50, 30, 104],
    )

    pdf.output("/usr/src/meganet/MegaNet_Technical_Report.pdf")
    print("PDF generated: MegaNet_Technical_Report.pdf")


# ---------------------------------------------------------------------------
# PPT
# ---------------------------------------------------------------------------
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BLUE   = RGBColor(0x0F, 0x1E, 0x3C)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD7)
LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xCC, 0xD6, 0xE8)


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def text_box(slide, text: str, left, top, width, height,
             size=Pt(16), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = size
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def multi_text_box(slide, lines: list, left, top, width, height,
                   default_size=Pt(16), default_color=WHITE, default_bold=False):
    """Lines is a list of (text, size, color, bold) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, col, bld) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = txt
        run.font.size = sz
        run.font.color.rgb = col
        run.font.bold = bld
        run.font.name = "Calibri"
    return txBox


def code_box(slide, text: str, left, top, width, height):
    add_rect(slide, left, top, width, height, LIGHT_BLUE)
    txBox = slide.shapes.add_textbox(
        left + Cm(0.2), top + Cm(0.2),
        width - Cm(0.4), height - Cm(0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.name = "Courier New"
        run.font.color.rgb = DARK_BLUE


def slide_header(slide, title: str):
    add_rect(slide, Cm(0), Cm(0), Cm(33.87), Cm(2.2), DARK_BLUE)
    add_rect(slide, Cm(0), Cm(2.2), Cm(33.87), Cm(0.18), ACCENT_BLUE)
    text_box(slide, title, Cm(0.8), Cm(0.2), Cm(32), Cm(2.0),
             size=Pt(24), bold=True, color=WHITE)


def build_ppt():
    prs = Presentation()
    prs.slide_width  = Cm(33.87)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    # ---------- Slide 1: Cover ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    add_rect(s, Cm(0), Cm(17.5), Cm(33.87), Cm(1.55), ACCENT_BLUE)

    multi_text_box(s, [
        ("MegaNet MVP", Pt(52), WHITE, True),
        ("Internet Paralela Descentralizada", Pt(26), ACCENT_BLUE, False),
        ("Kim Dotcom's Vision  |  Python Prototype  |  2026-02-23", Pt(15), LIGHT_GREY, False),
    ], Cm(2), Cm(4.5), Cm(29), Cm(7))

    text_box(s, "Hacking Conference 2026  |  Security Research Track",
             Cm(1), Cm(17.6), Cm(32), Cm(1.3),
             size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ---------- Slide 2: Problem ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "The Problem: Internet is Centralised & Censored")
    lines = [
        ("The centralised internet has critical weaknesses:", Pt(17), ACCENT_BLUE, True),
        ("", Pt(12), WHITE, False),
        ("  IP addresses are assigned by IANA/RIRs  ->  can be revoked or blocked", Pt(15), WHITE, False),
        ("  DNS is centralised  ->  single point of control and failure", Pt(15), WHITE, False),
        ("  BGP routing is trust-based  ->  route hijacking, censorship", Pt(15), WHITE, False),
        ("  Traffic is monitorable at ISP level", Pt(15), WHITE, False),
        ("", Pt(12), WHITE, False),
        ("Kim Dotcom's MegaNet vision:", Pt(17), ACCENT_BLUE, True),
        ("", Pt(12), WHITE, False),
        ("  No IP addresses  ->  addresses derived from public keys", Pt(15), WHITE, False),
        ("  No DNS  ->  blockchain-based key registry", Pt(15), WHITE, False),
        ("  Radio transport (LoRa)  ->  censorship-resistant physical layer", Pt(15), WHITE, False),
        ("  Encrypted by default  ->  E2E, no plaintext on the wire", Pt(15), WHITE, False),
    ]
    multi_text_box(s, lines, Cm(1), Cm(2.6), Cm(31), Cm(16))

    # ---------- Slide 3: Architecture ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "Architecture: 5 Independent Layers")
    code_box(s,
        "Layer 1  CRYPTO        Ed25519 identity + X25519 ECDH + ChaCha20-Poly1305\n"
        "Layer 2  BLOCKCHAIN    SHA3-256 PoW  |  4 TX types  |  fork resolution\n"
        "Layer 3  LORA RADIO    EU868 propagation model  |  duty-cycle tracker\n"
        "Layer 4  KADEMLIA DHT  160-bit XOR routing  |  k=20  |  content-addressing\n"
        "Layer 5  NODE          Integration: send_message / receive_packets / mine_block\n"
        "\n"
        "Single external dependency: cryptography>=42.0.0\n"
        "42 unit tests | 0.16s runtime | no hardware required",
        Cm(1), Cm(2.8), Cm(31), Cm(7.5)
    )
    multi_text_box(s, [
        ("Each layer is independently testable and replaceable.", Pt(16), ACCENT_BLUE, True),
        ("Real LoRa hardware can replace the simulator with zero changes to upper layers.", Pt(15), WHITE, False),
    ], Cm(1), Cm(11.5), Cm(31), Cm(4))

    # ---------- Slide 4: Crypto ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "Cryptography: Public-Key Addressing & E2E Encryption")
    code_box(s,
        "# Node identity & address\n"
        "ed_priv, ed_pub = Ed25519.generate()          # signing / identity\n"
        "x_priv,  x_pub  = X25519.generate()           # ECDH key exchange\n"
        "address          = SHA3_256(ed_pub_raw)[:20]  # 160-bit, no registry needed\n"
        "\n"
        "# Message encryption A -> B\n"
        "shared = X25519(A.x_priv, B.x_pub)            # 32-byte secret\n"
        "aad    = A.addr + B.addr                       # authenticated metadata\n"
        "nonce  = os.urandom(12)                        # 96-bit random\n"
        "ct     = ChaCha20Poly1305.encrypt(shared, plaintext, aad, nonce)\n"
        "wire   = A.addr(20B) + nonce(12B) + ct",
        Cm(1), Cm(2.8), Cm(31), Cm(9)
    )
    multi_text_box(s, [
        ("ChaCha20-Poly1305: constant-time, hardware-independent, 256-bit security", Pt(15), ACCENT_BLUE, False),
        ("Perfect for IoT/LoRa: no AES hardware required  |  lower power consumption", Pt(15), WHITE, False),
    ], Cm(1), Cm(12.8), Cm(31), Cm(4))

    # ---------- Slide 5: Blockchain ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "Lightweight Blockchain: Registration, Routing & Receipts")
    code_box(s,
        "# 4 Transaction types (all Ed25519 signed, SHA3-256 tx_id)\n"
        "NODE_REGISTER    ->  publish ed25519_pub + x25519_pub\n"
        "ROUTING_UPDATE   ->  announce up to 10 known peers\n"
        "MESSAGE_RECEIPT  ->  non-repudiable delivery proof (msg_id + receiver)\n"
        "DATA_ANCHOR      ->  pin content_hash + frag_count on-chain\n"
        "\n"
        "# Proof of Work (demo settings)\n"
        "difficulty = 2   # 2 leading hex zeros = ~256 iterations = ~1ms/block\n"
        "hash = SHA3_256(JSON(block_header))  must start with '00'\n"
        "\n"
        "# Mempool: max 50 txs, FIFO eviction, dedup by tx_id\n"
        "# Fork resolution: longest valid chain wins",
        Cm(1), Cm(2.8), Cm(31), Cm(10)
    )
    multi_text_box(s, [
        ("State indices rebuilt from chain: node_registry, message_receipts, data_anchors", Pt(15), WHITE, False),
    ], Cm(1), Cm(14.2), Cm(31), Cm(2))

    # ---------- Slide 6: LoRa ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "LoRaWAN EU868: Radio Transport Without IP")
    code_box(s,
        "# Log-distance path loss model\n"
        "PL(d) = 112 dB + 10 * 2.7 * log10(d_km)    # n=2.7 rural/suburban\n"
        "RSSI  = 14 dBm - PL(d)                      # TX power 14 dBm (EU868 max)\n"
        "\n"
        "# Delivery probability (sigmoid)\n"
        "P(rx) = 1 / (1 + exp(-0.5 * (RSSI - sensitivity[SF])))\n"
        "\n"
        "# EU868 Duty cycle: 1% per hour (sliding window)\n"
        "budget = 36,000 ms / hour    # enforced per gateway\n"
        "\n"
        "RSSI @ 1km: -98 dBm   @ 5km: -117 dBm   @ 10km: -125 dBm   @ 15km: -130 dBm\n"
        "SF7 sensitivity: -123 dBm   ->  effective range ~7-9 km",
        Cm(1), Cm(2.8), Cm(31), Cm(9.5)
    )
    multi_text_box(s, [
        ("5 gateways on 20x20 km grid  |  Fragment size: 203 B (SF7: 222B - 19B header)", Pt(15), WHITE, False),
        ("Real LoRa hardware plug-in: replace LoRaSimulator.transmit() with SX1276 driver", Pt(15), ACCENT_BLUE, False),
    ], Cm(1), Cm(13.5), Cm(31), Cm(4))

    # ---------- Slide 7: DHT ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "Kademlia DHT & 'Homeless Data'")
    code_box(s,
        "# 160-bit XOR routing (matches 20-byte address space)\n"
        "bucket_idx = (our_addr XOR their_addr).bit_length() - 1   # 0..159\n"
        "k = 20   # contacts per bucket, LRS eviction\n"
        "\n"
        "# Content addressing (homeless data)\n"
        "content_hash = SHA3_256(data)[:20]           # 20-byte content address\n"
        "store_key    = SHA3_256(content_hash + frag_idx)[:20]\n"
        "\n"
        "# Fragmentation\n"
        "fragment_message(data, max=203B) -> (content_hash, [Fragment...])\n"
        "reassemble_message(fragments)    -> bytes | None   # None if incomplete",
        Cm(1), Cm(2.8), Cm(31), Cm(9)
    )
    multi_text_box(s, [
        ("Data has no fixed home  ->  resistant to node takedowns", Pt(15), ACCENT_BLUE, False),
        ("Any node holding fragments can forward/serve them", Pt(15), WHITE, False),
        ("Content hash verified on reassembly  ->  data integrity guaranteed", Pt(15), WHITE, False),
    ], Cm(1), Cm(13.0), Cm(31), Cm(4))

    # ---------- Slide 8: Message Flow ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "End-to-End Message Flow: Alice -> Bob")
    code_box(s,
        "1. Lookup Bob.x25519_pub in blockchain.node_registry\n"
        "2. shared_secret = X25519(Alice.x_priv, Bob.x_pub)\n"
        "3. wire = Alice.addr(20B) + nonce(12B) + ChaCha20Poly1305(plaintext)\n"
        "4. fragment_message(wire) -> chunks of <=203 bytes\n"
        "5. for chunk in chunks:\n"
        "     pkt = LoRaPacket(SF7, chunk)\n"
        "     LoRaSimulator.transmit(pkt, Alice.x, Alice.y)   # path loss + random\n"
        "6. Gateways receive -> deduplicate -> deliver to Bob\n"
        "7. Bob.receive_packets() -> reassemble -> decrypt -> Bob.inbox\n"
        "8. Bob creates MESSAGE_RECEIPT tx -> mine_block() -> broadcast\n"
        "9. On-chain proof: blockchain.message_receipts[msg_id]",
        Cm(1), Cm(2.8), Cm(31), Cm(11)
    )
    multi_text_box(s, [
        ("Partitioned nodes cannot receive  |  All state verifiable on-chain", Pt(15), WHITE, False),
    ], Cm(1), Cm(14.8), Cm(31), Cm(2))

    # ---------- Slide 9: Security ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "Security Properties & Threat Model")
    lines = [
        ("Threats Addressed:", Pt(17), ACCENT_BLUE, True),
        ("", Pt(12), WHITE, False),
        ("  Eavesdropping         ->  ChaCha20-Poly1305 E2E encryption", Pt(15), WHITE, False),
        ("  Impersonation         ->  Ed25519 signatures on all transactions", Pt(15), WHITE, False),
        ("  Replay attacks        ->  Random 96-bit nonce per message", Pt(15), WHITE, False),
        ("  Payload tampering     ->  Poly1305 16-byte authentication tag", Pt(15), WHITE, False),
        ("  IP-based censorship   ->  No IP addresses in the protocol", Pt(15), WHITE, False),
        ("  Sybil (partial)       ->  PoW registration cost", Pt(15), WHITE, False),
        ("", Pt(12), WHITE, False),
        ("MVP Limitations (known):", Pt(17), ACCENT_BLUE, True),
        ("", Pt(12), WHITE, False),
        ("  No key revocation  |  PoW difficulty=2 trivial (demo only)", Pt(15), WHITE, False),
        ("  No incentive model for relaying  |  No store-and-forward for offline nodes", Pt(15), WHITE, False),
    ]
    multi_text_box(s, lines, Cm(1), Cm(2.6), Cm(31), Cm(16))

    # ---------- Slide 10: Test Results ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "Test Results: 42/42 Tests Pass")
    code_box(s,
        "$ .venv/bin/pytest tests/ -v\n"
        "\n"
        "test_blockchain.py::test_genesis_deterministic              PASSED\n"
        "test_blockchain.py::test_fork_resolution_longer_chain_wins  PASSED\n"
        "test_blockchain.py::test_mempool_fifo_eviction              PASSED\n"
        "test_crypto.py::test_ecdh_symmetric                         PASSED\n"
        "test_crypto.py::test_tamper_detection_raises                PASSED\n"
        "test_lora.py::test_rssi_at_5km                              PASSED\n"
        "test_lora.py::test_duty_cycle_blocks_over_budget            PASSED\n"
        "test_lora.py::test_simulator_far_node_low_delivery          PASSED\n"
        "...\n"
        "42 passed in 0.16s",
        Cm(1), Cm(2.8), Cm(31), Cm(9.5)
    )
    multi_text_box(s, [
        ("No hardware  |  No network  |  Single dependency: cryptography>=42", Pt(15), WHITE, False),
        ("python3 -m venv .venv && .venv/bin/pip install -r requirements.txt", Pt(14), ACCENT_BLUE, False),
    ], Cm(1), Cm(13.5), Cm(31), Cm(3))

    # ---------- Slide 11: Demo ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "Live Demo: python main.py")
    code_box(s,
        "$ python main.py\n"
        "\n"
        "1. Node Creation\n"
        "   Alice: 4ae58a4f3022fe697a88d19edfa5b18350f1498c\n"
        "   Bob:   67005aedfde22015fe044a1811f0300638b9097f\n"
        "   [PASS] All addresses are 40 hex chars & unique\n"
        "\n"
        "2. Bootstrap  ->  [PASS] node_registry has 5 entries  |  height >= 1\n"
        "3. Alice->Bob ->  [PASS] Bob received correct plaintext\n"
        "4. Receipt    ->  [PASS] Chain grew to >= 2  |  message_receipts populated\n"
        "5. Partition  ->  [PASS] Partitioned node (Eve) received nothing\n"
        "6. Routing    ->  [PASS] Alice has >= 4 contacts\n"
        "7. DHT        ->  [PASS] 500B -> 3 fragments, reassembly correct\n"
        "8. LoRa       ->  [PASS] RSSI @ 5km = -116.9 dBm (expected ~-117 dBm)",
        Cm(1), Cm(2.8), Cm(31), Cm(11.5)
    )

    # ---------- Slide 12: Roadmap ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    slide_header(s, "Roadmap: From MVP to Production")
    lines = [
        ("Phase 1 (MVP - DONE): Pure Python simulation  |  42 passing tests  |  full E2E demo", Pt(15), ACCENT_BLUE, True),
        ("", Pt(12), WHITE, False),
        ("Phase 2 - Real Hardware:", Pt(15), ACCENT_BLUE, True),
        ("  Swap LoRaSimulator -> SX1276/SX1278 driver  |  Raspberry Pi / ESP32 nodes", Pt(15), WHITE, False),
        ("", Pt(12), WHITE, False),
        ("Phase 3 - Protocol Hardening:", Pt(15), ACCENT_BLUE, True),
        ("  Key revocation  |  PoW difficulty 6+  |  Difficulty adjustment  |  Store-and-forward", Pt(15), WHITE, False),
        ("", Pt(12), WHITE, False),
        ("Phase 4 - Incentive Layer:", Pt(15), ACCENT_BLUE, True),
        ("  Token reward for fragment storage & relay  |  Proof-of-bandwidth consensus", Pt(15), WHITE, False),
        ("", Pt(12), WHITE, False),
        ("Phase 5 - Application Layer:", Pt(15), ACCENT_BLUE, True),
        ("  MegaNet DNS equivalent  |  Browser extension  |  Mobile app", Pt(15), WHITE, False),
    ]
    multi_text_box(s, lines, Cm(1), Cm(2.6), Cm(31), Cm(16))

    # ---------- Slide 13: Closing ----------
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BLUE)
    add_rect(s, Cm(0), Cm(17.5), Cm(33.87), Cm(1.55), ACCENT_BLUE)

    multi_text_box(s, [
        ("The internet can be rebuilt.", Pt(40), WHITE, True),
        ("No IPs. No DNS. No censorship.", Pt(24), ACCENT_BLUE, False),
        ("", Pt(14), WHITE, False),
        ("Code: /usr/src/meganet  |  Run: python main.py", Pt(14), LIGHT_GREY, False),
    ], Cm(2), Cm(5), Cm(29), Cm(8))

    text_box(s, "MegaNet MVP  |  Hacking Conference 2026  |  Questions?",
             Cm(1), Cm(17.6), Cm(32), Cm(1.3),
             size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    prs.save("/usr/src/meganet/MegaNet_Conference.pptx")
    print("PPT generated: MegaNet_Conference.pptx")


if __name__ == "__main__":
    build_pdf()
    build_ppt()
    print("\nAll documents generated successfully.")
